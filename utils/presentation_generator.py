import os
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage, PageBreak, Frame, PageTemplate
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from reportlab.pdfgen import canvas
from datetime import datetime
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup

class PresentationGenerator:
    """Generate eye-catching technical presentations - 1 page per item"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.setup_custom_styles()
    
    def setup_custom_styles(self):
        """Setup custom paragraph styles for presentations"""
        self.title_style = ParagraphStyle(
            'PresentationTitle',
            parent=self.styles['Heading1'],
            fontSize=28,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=20,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        self.item_title_style = ParagraphStyle(
            'ItemTitle',
            parent=self.styles['Heading1'],
            fontSize=22,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=15,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        self.spec_heading_style = ParagraphStyle(
            'SpecHeading',
            fontSize=14,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=8,
            fontName='Helvetica-Bold'
        )
        
        self.spec_text_style = ParagraphStyle(
            'SpecText',
            fontSize=11,
            textColor=colors.black,
            spaceAfter=6,
            alignment=TA_JUSTIFY
        )

    def _get_logo_path(self):
        candidates = [
            os.path.join('static', 'images', 'AlShaya-Logo-color@2x.png'),
            os.path.join('static', 'images', 'LOGO.png'),
            os.path.join('static', 'images', 'al-shaya-logo-white@2x.png')
        ]
        for p in candidates:
            if os.path.exists(p):
                return p
        return None

    def _draw_header_footer(self, canv: canvas.Canvas, doc):
        page_width, page_height = doc.pagesize
        gold = colors.HexColor('#d4af37')
        dark = colors.HexColor('#1a365d')
        # header rule
        canv.setStrokeColor(gold)
        canv.setLineWidth(2)
        canv.line(doc.leftMargin, page_height - doc.topMargin + 10, page_width - doc.rightMargin, page_height - doc.topMargin + 10)
        # logo top-right
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                w, h = 90, 90
                x = page_width - doc.rightMargin - w
                y = page_height - doc.topMargin + 12
                canv.drawImage(logo, x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
            except Exception:
                pass
        # footer site
        canv.setFillColor(dark)
        canv.setFont('Helvetica', 9)
        canv.drawRightString(page_width - doc.rightMargin, doc.bottomMargin - 12, 'https://alshayaenterprises.com')
        canv.setStrokeColor(gold)
        canv.setLineWidth(1)
        canv.line(doc.leftMargin, doc.bottomMargin - 8, page_width - doc.rightMargin, doc.bottomMargin - 8)
    
    def generate(self, file_id, session, format_type='pdf'):
        """
        Generate technical presentation with 1 page/slide per item
        Args:
            file_id: The file ID
            session: Flask session
            format_type: 'pdf' or 'pptx'
        Returns: path to generated file
        """
        # Get file info and extracted data
        uploaded_files = session.get('uploaded_files', [])
        file_info = None
        
        for f in uploaded_files:
            if f['id'] == file_id:
                file_info = f
                break
        
        if not file_info:
            raise Exception('File not found. Please upload and extract a file first.')
        
        # Get costed data (preferred) or stitched table or extraction result
        if 'costed_data' in file_info:
            items = self.parse_items_from_costed_data(file_info['costed_data'], session, file_id)
        elif 'stitched_table' in file_info:
            items = self.parse_items_from_stitched_table(file_info['stitched_table'], session, file_id)
        elif 'extraction_result' in file_info:
            items = self.parse_items_from_extraction(file_info['extraction_result'], session, file_id)
        else:
            raise Exception('No data available. Please extract tables first.')
        
        if not items:
            raise Exception('No items found in the table.')
        
        # Create output directory
        session_id = session['session_id']
        output_dir = os.path.join('outputs', session_id, 'presentations')
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate file based on format
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if format_type == 'pptx':
            output_file = os.path.join(output_dir, f'presentation_{file_id}_{timestamp}.pptx')
            self.generate_pptx(items, output_file)
        else:  # pdf
            output_file = os.path.join(output_dir, f'presentation_{file_id}_{timestamp}.pdf')
            self.generate_pdf(items, output_file)
        
        return output_file
    
    def parse_items_from_costed_data(self, costed_data, session, file_id):
        """Parse items from costed table data"""
        items = []
        session_id = session.get('session_id', '')
        
        for table in costed_data.get('tables', []):
            headers = [h for h in table.get('headers', []) if h.lower() not in ['action', 'actions']]
            
            for row in table.get('rows', []):
                # Find description column
                description = ''
                for h in headers:
                    if 'description' in h.lower() or 'item' in h.lower():
                        description = self.strip_html(row.get(h, ''))
                        break
                
                # Find quantity
                qty = ''
                unit = ''
                for h in headers:
                    if 'qty' in h.lower() or 'quantity' in h.lower():
                        qty = self.strip_html(row.get(h, ''))
                    if 'unit' in h.lower() and 'rate' not in h.lower():
                        unit = self.strip_html(row.get(h, ''))
                
                # Find pricing
                unit_rate = ''
                total = ''
                for h in headers:
                    if 'rate' in h.lower() or 'price' in h.lower():
                        unit_rate = self.strip_html(row.get(h, ''))
                    if 'total' in h.lower() or 'amount' in h.lower():
                        total = self.strip_html(row.get(h, ''))
                
                # Find image
                image_path = None
                for h in headers:
                    cell_value = row.get(h, '')
                    if self.contains_image(cell_value):
                        image_path = self.extract_image_path(cell_value, session_id, file_id)
                        break
                
                if description:  # Only add if we have a description
                    item = {
                        'description': description,
                        'qty': qty,
                        'unit': unit,
                        'unit_rate': unit_rate,
                        'total': total,
                        'image_path': image_path,
                        'brand': self.extract_brand(description),
                        'specifications': self.extract_specifications(description)
                    }
                    items.append(item)
        
        return items
    
    def parse_items_from_stitched_table(self, stitched_table, session, file_id):
        """Parse items from stitched HTML table data"""
        items = []
        session_id = session.get('session_id', '')
        
        # Parse the HTML
        html_content = stitched_table.get('html', '')
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Find the table
        table = soup.find('table')
        if not table:
            return items
        
        # Get headers
        headers = []
        header_row = table.find('tr')
        if header_row:
            for th in header_row.find_all(['th', 'td']):
                headers.append(th.get_text(strip=True).lower())
        
        # Get data rows (skip header row)
        rows = table.find_all('tr')[1:]  # Skip first row (headers)
        
        for row in rows:
            cells = row.find_all('td')
            if len(cells) != len(headers):
                continue
            
            # Build row dict
            row_data = {}
            for i, cell in enumerate(cells):
                if i < len(headers):
                    # Check if cell contains image
                    img = cell.find('img')
                    if img:
                        row_data[headers[i]] = str(cell)  # Keep HTML with image
                    else:
                        row_data[headers[i]] = cell.get_text(strip=True)
            
            # Extract fields
            description = ''
            for h in headers:
                if 'description' in h or 'item' in h:
                    description = self.strip_html(row_data.get(h, ''))
                    break
            
            if not description:
                continue
            
            # Find quantity
            qty = ''
            unit = ''
            for h in headers:
                if 'qty' in h or 'quantity' in h:
                    qty = self.strip_html(row_data.get(h, ''))
                if 'unit' in h and 'rate' not in h and 'price' not in h:
                    unit = self.strip_html(row_data.get(h, ''))
            
            # Find pricing
            unit_rate = ''
            total = ''
            for h in headers:
                if ('rate' in h or 'price' in h) and 'unit' in h:
                    unit_rate = self.strip_html(row_data.get(h, ''))
                if 'total' in h or 'amount' in h:
                    total = self.strip_html(row_data.get(h, ''))
            
            # Find image
            image_path = None
            for h in headers:
                cell_value = row_data.get(h, '')
                if self.contains_image(str(cell_value)):
                    image_path = self.extract_image_path(str(cell_value), session_id, file_id)
                    break
            
            item = {
                'description': description,
                'qty': qty,
                'unit': unit,
                'unit_rate': unit_rate,
                'total': total,
                'image_path': image_path,
                'brand': self.extract_brand(description),
                'specifications': self.extract_specifications(description)
            }
            items.append(item)
        
        return items
    
    def strip_html(self, text):
        """Strip HTML tags from text"""
        return re.sub(r'<[^>]+>', '', str(text)).strip()
    
    def contains_image(self, cell_value):
        """Check if cell contains an image reference"""
        return '<img' in str(cell_value).lower() or 'img_in_' in str(cell_value).lower()
    
    def extract_image_path(self, cell_value, session_id, file_id):
        """Extract image path from cell value"""
        try:
            match = re.search(r'src=["\']([^"\']+)["\']', str(cell_value))
            if match:
                img_relative_path = match.group(1).lstrip('/')
                if img_relative_path.startswith('outputs'):
                    return img_relative_path
                else:
                    return os.path.join('outputs', session_id, file_id, img_relative_path)
            
            if 'img_in_' in str(cell_value):
                match = re.search(r'(imgs/img_in_[^"\s<>]+\.jpg)', str(cell_value))
                if match:
                    img_relative_path = match.group(1)
                    return os.path.join('outputs', session_id, file_id, img_relative_path)
        except Exception as e:
            pass
        return None
    
    def generate_pdf(self, items, output_file):
        """Generate PDF presentation"""
        doc = SimpleDocTemplate(output_file, pagesize=A4, 
                                topMargin=0.5*inch, bottomMargin=0.5*inch,
                                leftMargin=0.75*inch, rightMargin=0.75*inch)
        story = []
        
        # Cover page
        story.extend(self.create_cover_page())
        story.append(PageBreak())
        
        # Create one page per item
        for idx, item in enumerate(items):
            story.extend(self.create_item_page_pdf(item, idx + 1))
            if idx < len(items) - 1:
                story.append(PageBreak())
        
        # Build PDF with header/footer
        doc.build(story, onFirstPage=self._draw_header_footer, onLaterPages=self._draw_header_footer)
    
    def generate_pptx(self, items, output_file):
        """Generate PowerPoint presentation"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        self.create_title_slide_pptx(prs)
        
        # Add one slide per item
        for idx, item in enumerate(items):
            self.create_item_slide_pptx(prs, item, idx + 1)
        
        prs.save(output_file)
    
    def create_title_slide_pptx(self, prs):
        """Create PowerPoint title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Big centered logo on title slide
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, Inches(3), Inches(0.6), width=Inches(4))
            except Exception:
                pass
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "TECHNICAL PROPOSAL"
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(26, 54, 93)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(3.7), Inches(6), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Furniture, Fixtures & Equipment"
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = RGBColor(26, 54, 93)
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(4), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime('%B %d, %Y')
        date_p = date_frame.paragraphs[0]
        date_p.font.size = Pt(16)
        date_p.alignment = PP_ALIGN.CENTER
    
    def create_item_slide_pptx(self, prs, item, page_num):
        """Create PowerPoint slide for one item"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Small logo top-right on content slides
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, Inches(8.4), Inches(0.2), width=Inches(1.4))
            except Exception:
                pass
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_text = f"Item {page_num}: {item['description'][:70]}"
        title_frame.text = title_text
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(20)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(26, 54, 93)
        
        # Image (left side)
        if item.get('image_path') and os.path.exists(item['image_path']):
            try:
                slide.shapes.add_picture(item['image_path'], Inches(0.5), Inches(1.2), 
                                        width=Inches(4), height=Inches(4))
            except Exception as e:
                # If image fails, add placeholder
                img_placeholder = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4), Inches(1))
                img_frame = img_placeholder.text_frame
                img_frame.text = "[Image Not Available]"
                img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        else:
            # Add placeholder
            img_placeholder = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4), Inches(1))
            img_frame = img_placeholder.text_frame
            img_frame.text = "[Image Not Available]"
            img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Details (right side)
        details_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(5.5))
        details_frame = details_box.text_frame
        details_frame.word_wrap = True
        
        # Product Details heading
        p = details_frame.paragraphs[0]
        p.text = "Product Details"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 54, 93)
        p.space_after = Pt(10)
        
        # Add details
        details_text = f"""
Brand: {item['brand']}
Quantity: {item['qty']} {item['unit']}
Unit Rate: {item['unit_rate']}
Total Amount: {item['total']}

Specifications:
"""
        for line in details_text.strip().split('\n'):
            if line.strip():
                p = details_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.space_after = Pt(6)
        
        # Add specifications
        for spec in item['specifications'][:5]:  # Limit to 5 specs
            p = details_frame.add_paragraph()
            p.text = f"• {spec}"
            p.font.size = Pt(12)
            p.level = 1
            p.space_after = Pt(4)
    
    def create_item_page_pdf(self, item, page_num):
        """Create PDF page for one item"""
        story = []
        
        # Item title
        item_title = f"Item {page_num}: {item['description'][:80]}"
        story.append(Paragraph(item_title, self.item_title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # Create two-column layout
        left_content = []
        right_content = []
        
        # Left: Image
        if item.get('image_path') and os.path.exists(item['image_path']):
            try:
                img = RLImage(item['image_path'], width=2.5*inch, height=2.5*inch)
                left_content.append(img)
            except Exception as e:
                left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        else:
            left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        
        # Right: Details
        details_html = f"""
            <para>
                <b><font size="14" color="#1a365d">Product Details</font></b><br/>
                <br/>
                <b>Brand:</b> {item['brand']}<br/>
                <b>Quantity:</b> {item['qty']} {item['unit']}<br/>
                <b>Unit Rate:</b> {item['unit_rate']}<br/>
                <b>Total Amount:</b> {item['total']}<br/>
                <br/>
                <b><font color="#1a365d">Specifications:</font></b><br/>
            </para>
        """
        right_content.append(Paragraph(details_html, self.spec_text_style))
        
        # Specifications
        for spec in item['specifications']:
            right_content.append(Paragraph(f"• {spec}", self.spec_text_style))
            right_content.append(Spacer(1, 0.05*inch))
        
        # Two-column table
        data = [[left_content, right_content]]
        t = Table(data, colWidths=[3*inch, 3.5*inch])
        t.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 10),
            ('RIGHTPADDING', (0, 0), (-1, -1), 10),
        ]))
        
        story.append(t)
        
        return story
    
    def parse_items_from_extraction(self, extraction_result, session, file_id):
        """
        Parse individual items from extraction result
        Returns: list of item dictionaries
        """
        items = []
        
        for layout_result in extraction_result.get('layoutParsingResults', []):
            markdown_text = layout_result.get('markdown', {}).get('text', '')
            images = layout_result.get('markdown', {}).get('images', {})
            
            # Parse tables from markdown
            table_rows = self.extract_table_rows(markdown_text)
            
            for row in table_rows:
                item = {
                    'sn': row.get('sn', row.get('sl.no', '')),
                    'description': row.get('description', row.get('item', '')),
                    'qty': row.get('qty', row.get('quantity', '')),
                    'unit': row.get('unit', ''),
                    'unit_rate': row.get('unit rate', row.get('unit price', row.get('rate', ''))),
                    'total': row.get('total', row.get('amount', '')),
                    'image': self.find_item_image(row, images),
                    'brand': self.extract_brand(row.get('description', '')),
                    'specifications': self.extract_specifications(row.get('description', ''))
                }
                items.append(item)
        
        return items
    
    def extract_table_rows(self, markdown_text):
        """Extract table rows from markdown text"""
        lines = markdown_text.split('\n')
        rows = []
        headers = []
        
        for line in lines:
            if '|' in line:
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                
                if not headers and cells:
                    # First row is headers
                    headers = [h.lower() for h in cells]
                elif headers and cells and not all(c in '-: ' for c in ''.join(cells)):
                    # Data row
                    if len(cells) == len(headers):
                        row = dict(zip(headers, cells))
                        rows.append(row)
        
        return rows
    
    def find_item_image(self, row, images):
        """Find image associated with this item"""
        # Try to find image reference in row
        for key, value in row.items():
            if 'image' in key.lower():
                # Check if value references an image
                for img_path, img_url in images.items():
                    if value in img_path or img_path in value:
                        return img_url
        
        # Return first available image if no specific match
        if images:
            return list(images.values())[0]
        
        return None
    
    def extract_brand(self, description):
        """Extract brand name from description (simple heuristic)"""
        # Common brand patterns - this is simplified
        brands = ['Sedus', 'Narbutas', 'Sokoa', 'B&T', 'Herman Miller', 'Steelcase', 'Haworth', 'Knoll']
        
        for brand in brands:
            if brand.lower() in description.lower():
                return brand
        
        # Try to extract first capitalized word
        words = description.split()
        for word in words:
            if word and word[0].isupper() and len(word) > 2:
                return word
        
        return 'Premium Brand'
    
    def extract_specifications(self, description):
        """Extract specifications from description"""
        # Split description into bullet points
        specs = []
        
        # Look for dimensions
        dimension_pattern = r'\d+\s*[xX×]\s*\d+\s*[xX×]?\s*\d*\s*(mm|cm|m|inch|in|")'
        dimensions = re.findall(dimension_pattern, description)
        if dimensions:
            specs.append(f"Dimensions: {', '.join(dimensions)}")
        
        # Look for materials
        materials = ['wood', 'metal', 'steel', 'aluminum', 'fabric', 'leather', 'plastic', 'glass', 'laminate']
        found_materials = [mat for mat in materials if mat in description.lower()]
        if found_materials:
            specs.append(f"Materials: {', '.join(found_materials).title()}")
        
        # Look for colors
        colors_list = ['black', 'white', 'grey', 'gray', 'brown', 'blue', 'red', 'green', 'beige']
        found_colors = [col for col in colors_list if col in description.lower()]
        if found_colors:
            specs.append(f"Available Colors: {', '.join(found_colors).title()}")
        
        if not specs:
            # Use description as-is if no specific specs found
            specs.append(description[:200])
        
        return specs
    
    def create_cover_page(self):
        """Create presentation cover page"""
        story = []
        
        # Centered large logo (if available)
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                img = RLImage(logo, width=3.5*inch, height=3.5*inch)
                img.hAlign = 'CENTER'
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
            except Exception:
                pass

        # Title
        title = Paragraph("TECHNICAL PROPOSAL", self.title_style)
        story.append(Spacer(1, 2*inch))
        story.append(title)
        story.append(Spacer(1, 0.5*inch))
        
        # Subtitle
        subtitle_style = ParagraphStyle(
            'Subtitle',
            parent=self.styles['Normal'],
            fontSize=16,
            textColor=colors.HexColor('#1a365d'),
            alignment=TA_CENTER
        )
        subtitle = Paragraph("Furniture, Fixtures & Equipment", subtitle_style)
        story.append(subtitle)
        story.append(Spacer(1, 1*inch))
        
        # Company info
        company_info = f"""
            <para align="center">
                <b>Prepared By:</b><br/>
                <font size="14" color="#667eea"><b>Your Company Name</b></font><br/>
                <br/>
                Date: {datetime.now().strftime('%B %d, %Y')}<br/>
            </para>
        """
        story.append(Paragraph(company_info, self.styles['Normal']))
        
        return story
    
    def create_item_page(self, item, page_num):
        """Create one page for an item with eye-catching design"""
        story = []
        
        # Item number and title
        item_title = f"Item {page_num}: {item['description'][:60]}"
        story.append(Paragraph(item_title, self.item_title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Create two-column layout using table
        left_content = []
        right_content = []
        
        # Left column - Image
        if item['image']:
            try:
                # For now, placeholder - in production, download and embed image
                img_placeholder = Paragraph(
                    f'<para align="center"><b>[Product Image]</b><br/>{item["image"][:50]}...</para>',
                    self.styles['Normal']
                )
                left_content.append(img_placeholder)
            except:
                left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        else:
            left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        
        # Right column - Specifications
        specs_html = f"""
            <para>
                <b><font size="14" color="#667eea">Product Details</font></b><br/>
                <br/>
                <b>Brand:</b> {item['brand']}<br/>
                <b>Quantity:</b> {item['qty']} {item['unit']}<br/>
                <b>Unit Rate:</b> {item['unit_rate']}<br/>
                <b>Total Amount:</b> {item['total']}<br/>
                <br/>
                <b><font color="#667eea">Specifications:</font></b><br/>
            </para>
        """
        right_content.append(Paragraph(specs_html, self.spec_text_style))
        
        # Add specifications as bullet points
        for spec in item['specifications']:
            spec_bullet = f"• {spec}"
            right_content.append(Paragraph(spec_bullet, self.spec_text_style))
            right_content.append(Spacer(1, 0.1*inch))
        
        # Additional info
        additional_info = """
            <para>
                <br/>
                <b><font color="#667eea">Additional Information:</font></b><br/>
                • Country of Origin: Various<br/>
                • Warranty: As per manufacturer standard<br/>
                • Lead Time: 4-6 weeks<br/>
                • Finish: As specified or equivalent<br/>
            </para>
        """
        right_content.append(Paragraph(additional_info, self.spec_text_style))
        
        # Create two-column table
        data = [[left_content, right_content]]
        
        col_widths = [3*inch, 3.5*inch]
        t = Table(data, colWidths=col_widths)
        t.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 10),
            ('RIGHTPADDING', (0, 0), (-1, -1), 10),
        ]))
        
        story.append(t)
        story.append(Spacer(1, 0.3*inch))
        
        # Bottom section - Key features
        features_title = Paragraph('<b><font size="12" color="#1a365d">KEY FEATURES</font></b>', self.styles['Normal'])
        story.append(features_title)
        story.append(Spacer(1, 0.1*inch))
        
        features = [
            "✓ Premium quality construction",
            "✓ Modern ergonomic design",
            "✓ Environmentally friendly materials",
            "✓ Easy maintenance and durability"
        ]
        
        features_text = '<br/>'.join(features)
        story.append(Paragraph(features_text, self.spec_text_style))
        
        return story
